"""
File: mdmt_doctoppt.py
Author: 林星語
Create Date: 4.7.2024
Description: This Python script transfers a word file to ppt 
"""

import os
import docx
from pptx import Presentation


def set_context(slide, context, keyword):
    """
    set context to specific fields
 
    Args:
        slide (Slide): specific slide.
        context (string): words to fill in the slides
        keyword (string): 標題欄位/人欄位.        
 
    Returns:
        none
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if keyword in run.text:
                    run.text = context

def get_context_size(slide,  keyword):
    """
    calculate the number of words on a particular slide
 
    Args:
        slide (Slide): specific slide.
        keyword (string): 標題欄位/人欄位.        
 
    Returns:
        the number of words in each slide

    """    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if keyword in run.text:
                    return len(run.text)
    return 0

def find_class_title(str):
    """
    check if the class titles exist in the specific string

    Args:
        str (string): specific slide
 
    Returns:
        True:a class title exists in the line
        None:there isn't a class title in the line

    """
    for i in range(len(str)-2):
        if "國高美餐應廣英子商觀訊".find(str[i]) > -1:
            if "一二三際".find(str[i+1]) > -1:
                if "0123456789０１２３４５６７８９".find(str[i+2]) > -1:
                    return True
    return None

def split_line(context, num):
    """
    consider the breakpoints of the fields and calculate the length of each line. Divide them so they can fit the fields on the ppt if necessary
    
    Args:
        context (string): original paragraph
        num (int):limit width of each field
 

    Returns:
        result(list):the elements in the list are each line after segemtation
        
    """
    tk1 = context.split()
    tokens = []
    for t in tk1:
        # get rid of the token'*' in the lines
        t = t.replace("*", "")
        t = t.replace("。", "")
        # split the strings with the token'、'
        tokens.extend(t.split('、'))

    result = []
    result.append('')
    seg = ''
    for t in tokens:
        if(find_class_title(t)):
            # class
            if(len(seg)+1+len(result[-1]) < num):
                if len(result[-1]) > 0:
                    result[-1] = result[-1]+" "
                result[-1] = result[-1]+seg
                seg = t
            else:
                result.append(seg.strip())
                seg = t

        else:
            if len(seg) > 0:
                seg = seg + " "
            seg = seg+t
    if len(seg) > 0:
        if(len(seg)+1+len(result[-1]) < num):
            result[-1] = result[-1]+" "+seg
        else:
            result.append(seg.strip())

    fresult = []
    for line in result :
        fresult.append( line.strip())
    return fresult

def set_context_to_slide(prs, context, title, slide_idx):
    """
    set title context to each slide
    Args:
        prs (Presentation): ppt file obj
        context (string): string context
        title (string): slide title
        slide_idx (string): slide index

    Returns:
        slide index        
    """
    if prs.slides[slide_idx] != None and len(context) > 0:
        i = 0
        while(i < len(context)):
            if len(context[i].strip())==0 :
                i=i+1
                continue 
            # 超過四行->下一頁
            if(len(context)-i > 4):
                c = context[i]+'\n'+context[i+1]+'\n'+context[i+2]
                slide = prs.slides[slide_idx]
                set_context(slide, title, "標題欄位")
                set_context(slide, c, "人欄位")
                slide_idx = slide_idx+1
                i = i+3
            # 內容等於四行->分兩列
            elif len(context)-i == 4:
                c1 = context[i]+'\n'+context[i+1]
                c2 = context[i+2]+'\n'+context[i+3]
                slide = prs.slides[slide_idx]
                set_context(slide, title, "標題欄位")
                set_context(slide, c1, "人欄位")
                slide_idx = slide_idx+1
                slide = prs.slides[slide_idx]
                set_context(slide, title, "標題欄位")
                set_context(slide, c2, "人欄位")
                slide_idx = slide_idx+1
                i = i+4
            # 內容小於四行->不用換頁
            else:
                c = ''
                for k in range(i, len(context)):
                    if(k == i):
                        c = context[k]
                    else:
                        c = c+"\n"+context[k]
                slide = prs.slides[slide_idx]
                set_context(slide, title, "標題欄位")
                set_context(slide, c, "人欄位")
                slide_idx = slide_idx+1
                i = len(context)
    return slide_idx
               

def title_filter( title ):
    """
    find the "("s and ")"s in the titled and remove everthing instde
    
    Args:
        title (string): original title 

    Returns:
        title (string): the string without "()"
    """

    while( True ):
        a = title.find("(") 
        if( a==-1):
            a=title.find("（") 
        if( a>-1 ):                
            c = title.find(")",a) 
            if( c==-1):
                c=title.find("）",a) 
            if a>-1 and  c>-1 :
                title = title[0:a] + title[c+1:]
        else:
            break 
    return title

def parsing_doc(doc_file , output_file , template_file):

    doc = docx.Document(doc_file)
    context = []
    prs = Presentation(template_file)
    # 建立簡報檔第一張頁面物件
    slide_idx = 3
    title = ''

    slide = prs.slides[slide_idx]
    for para in doc.paragraphs:
        content = para.text.strip()
        if len(content) == 0:
            continue
        if content.find("*****") > -1:
            continue

        m = find_class_title(content)
        if m == None:  # 出現標題
            print("==>", content)

            

            slide_idx = set_context_to_slide(prs, context, title, slide_idx)
                
            title = title_filter(content)
            print("==>", title)
            context = []
        else:
            if len(content) > 0:
                context.extend(split_line(content, 35))

    slide_idx = set_context_to_slide(prs, context, title, slide_idx)

    # 將簡報物件存檔
    prs.save(output_file)
 
folder = "files"
for f in os.listdir(folder):
    if f.endswith(".docx"):
        # print(f.)
        parsing_doc(folder+'/'+f , folder+'/'+f+'.pptx' , 't.pptx')